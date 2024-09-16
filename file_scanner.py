import magic
import pdfplumber
from pptx import Presentation
import pefile
from utils import calculate_hash, speak
import logging

def scan_file(file_path, model, vectorizer, app):
    """Scan a file and make predictions based on the file's content."""
    mime = magic.Magic(mime=True)
    mime_type = mime.from_file(file_path)
    file_hash = calculate_hash(file_path)
    app.insert_text(f"File MIME type: {mime_type}, Hash: {file_hash}\n", "info")

    if mime_type in ["application/pdf", "application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/x-dosexec"]:
        text = extract_text(file_path, mime_type)
        if text:
            features = vectorizer.transform([text])
            prediction = model.predict(features)
            result = "Malicious" if prediction[0] else "Benign"
            app.insert_text(f"BenardAI: The file is {result}.\n", "result")
            speak(f"The file is {result}.")
        else:
            app.insert_text("BenardAI: No text extracted from the file.\n", "error")
            speak("No text extracted from the file.")
    else:
        app.insert_text("BenardAI: Unsupported file type.\n", "error")
        speak("Unsupported file type.")

def extract_text(file_path, mime_type):
    """Extract text from various file types."""
    try:
        if mime_type == "application/pdf":
            return extract_text_from_pdf(file_path)
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return extract_text_from_ppt(file_path)
        elif mime_type == "application/x-dosexec":
            return extract_info_from_exe(file_path)
    except Exception as e:
        logging.error(f"Error extracting text: {e}")
    return ""

def extract_text_from_pdf(file_path):
    """Extract text from a PDF file."""
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text

def extract_text_from_ppt(file_path):
    """Extract text from a PowerPoint presentation."""
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def extract_info_from_exe(file_path):
    """Extract metadata from a PE executable."""
    pe = pefile.PE(file_path)
    info = [f"File Info: {pe.FILE_HEADER.Machine}", f"TimeDateStamp: {pe.FILE_HEADER.TimeDateStamp}"]
    return "\n".join(info)
