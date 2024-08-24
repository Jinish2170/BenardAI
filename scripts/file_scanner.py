import magic
import hashlib
import logging
import pdfplumber
from pptx import Presentation
import pefile

class FileScanner:
    def __init__(self, model_manager):
        self.model_manager = model_manager

    def scan_file(self, file_path):
        mime = magic.Magic(mime=True)
        mime_type = mime.from_file(file_path)
        file_hash = self.calculate_hash(file_path)
        logging.info(f"File MIME type: {mime_type}, Hash: {file_hash}")

        if mime_type in ["application/pdf", "application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/x-dosexec"]:
            if mime_type == "application/pdf":
                text = self.extract_text_from_pdf(file_path)
            elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                text = self.extract_text_from_ppt(file_path)
            elif mime_type == "application/x-dosexec":
                text = self.extract_info_from_exe(file_path)
            logging.info(f"Extracted text: {text[:500]}...")  # Log only first 500 characters

            features = self.model_manager.vectorize_text(text)
            prediction = self.model_manager.predict(features)
            result = "Malicious" if prediction[0] else "Benign"

            logging.info(f"Scan result: {result}")
        else:
            logging.error("Unsupported file type.")

    def calculate_hash(self, file_path):
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()

    def extract_text_from_pdf(self, file_path):
        text = ""
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text += page.extract_text()
        except Exception as e:
            logging.error(f"Error extracting text from PDF: {e}")
        return text

    def extract_text_from_ppt(self, file_path):
        text = ""
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text
        except Exception as e:
            logging.error(f"Error extracting text from PPT: {e}")
        return text

    def extract_info_from_exe(self, file_path):
        text = ""
        try:
            pe = pefile.PE(file_path)
            for entry in pe.DIRECTORY_ENTRY_IMPORT:
                text += f"{entry.dll.decode()}\n"
        except Exception as e:
            logging.error(f"Error extracting info from EXE: {e}")
        return text
